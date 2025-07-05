"""
Document processing module for handling various file formats.
"""
import os
import tempfile
from typing import List, Optional, Union
from pathlib import Path
import numpy as np

# Document processing imports
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import openai
from loguru import logger

from config import get_settings

settings = get_settings()


class LightweightEmbeddingService:
    """
    Lightweight embedding service using OpenAI's text-embedding-3-small model.

    This service uses the text-embedding-3-small model which offers:
    - High performance with 61.0% MTEB score
    - Configurable dimensions (1536 max, can be reduced for cost savings)
    - Cost-effective at $0.00002 per 1K tokens
    - Multilingual support
    - Max input: 8,191 tokens

    Dimension Options:
    - 1536: Full quality (default)
    - 512: ~67% cost reduction, minimal quality loss
    - 256: ~83% cost reduction, moderate quality loss

    Environment Configuration:
    - EMBEDDING_MODEL: text-embedding-3-small (default)
    - EMBEDDING_DIMENSION: 1536 (default, configurable)
    """

    def __init__(self):
        """Initialize the embedding service."""
        self.client = openai.OpenAI(api_key=settings.openai_api_key)
        self.model = settings.embedding_model  # Configurable via environment
        self.dimension = settings.embedding_dimension  # Configurable via environment

    def encode(self, texts: List[str]) -> List[List[float]]:
        """
        Generate embeddings for texts using OpenAI API.

        Args:
            texts: List of text strings

        Returns:
            List of embedding vectors
        """
        try:
            # Handle batch processing for large text lists
            batch_size = 100  # OpenAI rate limits
            all_embeddings = []

            for i in range(0, len(texts), batch_size):
                batch = texts[i:i + batch_size]

                response = self.client.embeddings.create(
                    model=self.model,
                    input=batch,
                    dimensions=self.dimension  # Configurable dimension for cost optimization
                )

                batch_embeddings = [embedding.embedding for embedding in response.data]
                all_embeddings.extend(batch_embeddings)

            return all_embeddings

        except Exception as e:
            logger.error(f"Error generating embeddings: {str(e)}")
            # Fallback to simple vector representations
            return self._fallback_embeddings(texts)

    def _fallback_embeddings(self, texts: List[str]) -> List[List[float]]:
        """
        Fallback method for creating simple embeddings when API fails.
        Creates basic TF-IDF like vectors for demonstration.
        """
        logger.warning("Using fallback embedding method")

        # Simple word-based vectors for fallback
        embeddings = []
        for text in texts:
            words = text.lower().split()
            # Create a simple hash-based vector
            vector = [0.0] * 384  # Smaller dimension for fallback
            for word in words:
                hash_val = hash(word) % 384
                vector[hash_val] += 1.0

            # Normalize
            norm = sum(x * x for x in vector) ** 0.5
            if norm > 0:
                vector = [x / norm for x in vector]

            embeddings.append(vector)

        return embeddings


class DocumentProcessor:
    """Handles processing of various document formats."""

    def __init__(self):
        """Initialize the document processor."""
        self.embedding_service = LightweightEmbeddingService()
        self.supported_formats = {
            '.txt': self._process_text,
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx,
        }

    def process_file(self, file_path: Union[str, Path]) -> Optional[List[str]]:
        """
        Process a file and extract text content.

        Args:
            file_path: Path to the file to process

        Returns:
            List of text chunks or None if processing failed
        """
        file_path = Path(file_path)

        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            return None

        extension = file_path.suffix.lower()

        if extension not in self.supported_formats:
            logger.error(f"Unsupported file format: {extension}")
            return None

        try:
            processor = self.supported_formats[extension]
            text_content = processor(file_path)

            if text_content:
                chunks = self._chunk_text(text_content)
                logger.info(f"Successfully processed {file_path}, created {len(chunks)} chunks")
                return chunks
            else:
                logger.warning(f"No text content extracted from {file_path}")
                return None

        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            return None

    def _process_text(self, file_path: Path) -> str:
        """Process plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as file:
                return file.read()

    def _process_pdf(self, file_path: Path) -> str:
        """Process PDF file."""
        text_content = []

        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)

            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(f"[Page {page_num + 1}]\n{text}")
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num + 1}: {str(e)}")
                    continue

        return '\n\n'.join(text_content)

    def _process_docx(self, file_path: Path) -> str:
        """Process DOCX file."""
        doc = DocxDocument(str(file_path))
        text_content = []

        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(' | '.join(row_text))

        return '\n\n'.join(text_content)

    def _process_pptx(self, file_path: Path) -> str:
        """Process PPTX file."""
        presentation = Presentation(str(file_path))
        text_content = []

        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = [f"[Slide {slide_num}]"]

            # Extract text from shapes using dynamic attribute access
            for shape in slide.shapes:
                try:
                    # Use getattr to safely access attributes
                    if hasattr(shape, "text_frame"):
                        text_frame = getattr(shape, "text_frame", None)
                        if text_frame and hasattr(text_frame, "text"):
                            text = getattr(text_frame, "text", "").strip()
                            if text:
                                slide_text.append(text)
                    elif hasattr(shape, "text"):
                        text = getattr(shape, "text", "").strip()
                        if text:
                            slide_text.append(text)
                except Exception:
                    # Skip problematic shapes
                    pass

                # Extract table content
                try:
                    if hasattr(shape, "table"):
                        table = getattr(shape, "table", None)
                        if table and hasattr(table, "rows"):
                            for row in getattr(table, "rows", []):
                                row_text = []
                                if hasattr(row, "cells"):
                                    for cell in getattr(row, "cells", []):
                                        cell_text = getattr(cell, "text", "").strip()
                                        if cell_text:
                                            row_text.append(cell_text)
                                if row_text:
                                    slide_text.append(' | '.join(row_text))
                except Exception:
                    # Skip shapes without accessible tables
                    pass

            if len(slide_text) > 1:  # More than just the slide number
                text_content.append('\n'.join(slide_text))

        return '\n\n'.join(text_content)

    def _chunk_text(self, text: str) -> List[str]:
        """
        Split text into overlapping chunks.

        Args:
            text: Input text to chunk

        Returns:
            List of text chunks
        """
        if not text.strip():
            return []

        # Split by sentences/paragraphs for better semantic chunks
        sentences = text.replace('\n\n', ' [PARAGRAPH_BREAK] ').split('. ')

        chunks = []
        current_chunk = ""
        current_size = 0

        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue

            sentence_size = len(sentence.split())

            # If adding this sentence exceeds chunk size, save current chunk
            if current_size + sentence_size > settings.chunk_size and current_chunk:
                chunks.append(current_chunk.strip())

                # Start new chunk with overlap
                overlap_words = current_chunk.split()[-settings.chunk_overlap:]
                current_chunk = ' '.join(overlap_words) + ' ' + sentence
                current_size = len(overlap_words) + sentence_size
            else:
                current_chunk += ' ' + sentence
                current_size += sentence_size

        # Add the last chunk
        if current_chunk.strip():
            chunks.append(current_chunk.strip())

        # Clean up paragraph breaks
        cleaned_chunks = []
        for chunk in chunks:
            cleaned = chunk.replace(' [PARAGRAPH_BREAK] ', '\n\n')
            if cleaned.strip():
                cleaned_chunks.append(cleaned.strip())

        return cleaned_chunks

    def get_embeddings(self, texts: List[str]) -> List[List[float]]:
        """
        Generate embeddings for a list of texts.

        Args:
            texts: List of text strings

        Returns:
            List of embedding vectors
        """
        try:
            return self.embedding_service.encode(texts)
        except Exception as e:
            logger.error(f"Error generating embeddings: {str(e)}")
            return []


def save_uploaded_file(uploaded_file, upload_dir: str = "uploads") -> Optional[str]:
    """
    Save an uploaded file to the uploads directory.

    Args:
        uploaded_file: Streamlit uploaded file object
        upload_dir: Directory to save the file

    Returns:
        Path to the saved file or None if failed
    """
    try:
        os.makedirs(upload_dir, exist_ok=True)

        file_path = os.path.join(upload_dir, uploaded_file.name)

        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        logger.info(f"File saved: {file_path}")
        return file_path

    except Exception as e:
        logger.error(f"Error saving file: {str(e)}")
        return None
